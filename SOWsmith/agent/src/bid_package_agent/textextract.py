"""Extract plain text from the KB file types so the retriever can ground on all of them
(DOCX, PDF, PPTX, XLSX, TXT/MD) — the same content Copilot Studio's SharePoint knowledge
source indexes. PDF extraction is best-effort (pypdf, else the pdftotext CLI, else skip)."""
from __future__ import annotations

import shutil
import subprocess
from pathlib import Path


def _docx(path: Path) -> str:
    from docx import Document
    d = Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        return "\n".join((p.extract_text() or "") for p in PdfReader(str(path)).pages)
    except Exception:
        pass
    if shutil.which("pdftotext"):
        try:
            return subprocess.run(["pdftotext", "-q", str(path), "-"],
                                  capture_output=True, text=True, timeout=30).stdout
        except Exception:
            return ""
    return ""


def extract_text(path) -> str:
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext == ".docx":
            return _docx(p)
        if ext == ".pptx":
            return _pptx(p)
        if ext == ".xlsx":
            return _xlsx(p)
        if ext == ".pdf":
            return _pdf(p)
        if ext in (".txt", ".md"):
            return p.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""
    return ""
