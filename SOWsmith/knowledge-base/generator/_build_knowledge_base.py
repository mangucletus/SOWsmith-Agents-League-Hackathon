#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Builds the SharePoint Online knowledge base for the Bid Package Generator (SOW/RFP) POC for SOWsmith.

Output -> ../sharepoint-online-knowledge-base/
  Approved-Exemplars/   3-5 gold-standard prior SOWs   -> feeds the GENERATOR (Flow B)
  Reference-Library/    broad library of prior bid docs + clause library -> feeds Q&A (Flow A)
  Templates/            canonical SOW + RFP templates, clause library, bidder form
  Lists/                Reviewers.xlsx, AuditLog.xlsx (import as SharePoint lists)
  Drafts/ Approved/ Rejected/   output libraries (start empty)
  README.md             structure, scoping rules, upload instructions, manifest

Design rules (mirror the engine):
  * Grounding libraries contain ONLY bid documents (no readme inside them).
  * service_type tokens are byte-identical across footers and the Reviewers list.
  * Every SOW uses the SAME seven headings + footer format (the house style the generator copies).
  * Concrete, self-consistent values live in the exemplars (the generator's factual anchor).
  * US English throughout. Footer is BODY text (extracts cleanly for RAG).
  * Review status: STANDARD REVIEW, or LEGAL REVIEW REQUIRED when non-standard commercial/legal
    terms are present (liquidated damages, indemnification, bonding, etc.) -> Legal/Contracts.
"""

import os
import shutil
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PColor
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
_CANDIDATES = [os.path.join(HERE, os.pardir, "sharepoint-online-knowledge-base"),
               os.path.join(HERE, "sharepoint-online-knowledge-base")]
ROOT = os.path.abspath(next((c for c in _CANDIDATES if os.path.isdir(c)), _CANDIDATES[0]))

NAVY = RGBColor(0x0C, 0x3B, 0x4D)
AMBER = RGBColor(0xBD, 0x60, 0x15)
GREY = RGBColor(0x5B, 0x72, 0x7C)
ORG = "SOWsmith"
SEVEN = ["1. PROJECT OVERVIEW", "2. SCOPE OF WORK", "3. EXCLUSIONS & ASSUMPTIONS",
         "4. DELIVERABLES & SCHEDULE", "5. MATERIALS, EQUIPMENT & SITE CONDITIONS",
         "6. SAFETY, QUALITY & COMPLIANCE", "7. ACCEPTANCE & REFERENCES"]
STANDARD = "STANDARD REVIEW"
LEGAL = "LEGAL REVIEW REQUIRED"

# ---------------------------------------------------------------- docx helpers
def new_doc():
    doc = Document(); n = doc.styles["Normal"].font; n.name = "Calibri"; n.size = Pt(11); return doc

def _shade(cell, hexcolor):
    tcPr = cell._tc.get_or_add_tcPr(); shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto"); shd.set(qn("w:fill"), hexcolor)
    tcPr.append(shd)

def title_block(doc, title, subtitle):
    p = doc.add_paragraph(); r = p.add_run(title); r.bold = True; r.font.size = Pt(19); r.font.color.rgb = NAVY
    s = doc.add_paragraph(); rs = s.add_run(subtitle); rs.italic = True; rs.font.size = Pt(10.5); rs.font.color.rgb = GREY
    pr = doc.add_paragraph(); pPr = pr._p.get_or_add_pPr(); bdr = OxmlElement("w:pBdr"); bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single"); bottom.set(qn("w:sz"), "12"); bottom.set(qn("w:space"), "1"); bottom.set(qn("w:color"), "BD6015")
    bdr.append(bottom); pPr.append(bdr)

def control_table(doc, rows):
    t = doc.add_table(rows=0, cols=2); t.style = "Table Grid"; t.alignment = WD_TABLE_ALIGNMENT.LEFT
    for k, v in rows:
        c = t.add_row().cells; kr = c[0].paragraphs[0].add_run(k); kr.bold = True; kr.font.size = Pt(9.5); kr.font.color.rgb = NAVY
        _shade(c[0], "F1F5F7"); c[1].paragraphs[0].add_run(v).font.size = Pt(9.5)
    for row in t.rows:
        row.cells[0].width = Inches(1.8); row.cells[1].width = Inches(4.5)
    doc.add_paragraph()

def heading(doc, text):
    h = doc.add_heading(level=1); r = h.add_run(text); r.font.color.rgb = NAVY; r.font.size = Pt(13); r.bold = True

def para(doc, text):
    doc.add_paragraph(text)

def numbered(doc, items):
    for s in items:
        doc.add_paragraph(s, style="List Number")

def bullets(doc, items):
    for s in items:
        doc.add_paragraph(s, style="List Bullet")

def opt_table(doc, caption, headers, rows):
    cp = doc.add_paragraph(); cr = cp.add_run(caption); cr.bold = True; cr.font.size = Pt(10); cr.font.color.rgb = NAVY
    t = doc.add_table(rows=1, cols=len(headers)); t.style = "Table Grid"
    for i, h in enumerate(headers):
        _shade(t.rows[0].cells[i], "EAF0F3"); rr = t.rows[0].cells[i].paragraphs[0].add_run(h); rr.bold = True; rr.font.size = Pt(9)
    for row in rows:
        cells = t.add_row().cells
        for i, val in enumerate(row):
            cells[i].paragraphs[0].add_run(str(val)).font.size = Pt(9)
    doc.add_paragraph()

def footer_block(doc, version, date, service_type, document_type, status):
    doc.add_paragraph()
    for line in ["---", f"VERSION: {version}", f"DATE: {date}", f"SERVICE TYPE: {service_type}",
                 f"DOCUMENT TYPE: {document_type}", f"REVIEW STATUS: {status}", "---"]:
        p = doc.add_paragraph(); r = p.add_run(line); r.font.name = "Consolas"; r.font.size = Pt(9.5); r.font.color.rgb = GREY
        p.paragraph_format.space_after = Pt(0)

# ---------------------------------------------------------------- the builder
def build_sow(spec, out_path):
    doc = new_doc()
    title_block(doc, spec["title"], f"{ORG} · Bid Package · Supply Chain")
    control_table(doc, [
        ("Document ID", spec["docid"]), ("Service type", spec["type"]),
        ("Document type", spec.get("doc", "SOW")), ("Version", spec["version"]),
        ("Effective date", spec["date"]), ("Owner", "Supply Chain"),
        ("Classification", "Internal"), ("Review status", spec["status"]),
    ])
    heading(doc, SEVEN[0]); para(doc, spec["overview"])
    heading(doc, SEVEN[1]); numbered(doc, spec["scope"])
    heading(doc, SEVEN[2]); bullets(doc, spec["exclusions"])
    heading(doc, SEVEN[3]); para(doc, spec["deliverables"])
    if spec.get("schedule_table"):
        cap, hdrs, rows = spec["schedule_table"]; opt_table(doc, cap, hdrs, rows)
    heading(doc, SEVEN[4]); para(doc, spec["materials"])
    heading(doc, SEVEN[5]); para(doc, spec["safety"])
    heading(doc, SEVEN[6]); para(doc, spec["acceptance"])
    bullets(doc, spec["references"])
    footer_block(doc, spec["version"], spec["date"], spec["type"], spec.get("doc", "SOW"), spec["status"])
    doc.save(out_path)

# ================================================================ CONTENT
D = []  # list of bid-document specs
def add(**kw): D.append(kw)

_SAFETY = (f"All work is performed under {ORG}'s safety-first program. The contractor complies with "
           "OSHA construction standards, applicable DOT and PHMSA requirements, the site-specific "
           "Job Safety Analysis, and all permit-to-work controls. Personal protective equipment is "
           "mandatory. Quality follows the referenced specifications and inspection/test plan.")

# -------- EXEMPLAR 1: Pipeline Construction --------------------------------
add(type="Pipeline Construction", exemplar=True, status=STANDARD, doc="SOW",
 file="Pipeline-Construction-SOW", title=f"{ORG} Pipeline Construction — Statement of Work (Exemplar)",
 docid="BP-SOW-PLC-001", version="v2.1", date="2026-01-15",
 overview=("Construction of a buried steel natural-gas pipeline segment, including all civil, "
   "mechanical, welding and testing activities required to deliver a complete, tested, and "
   "ready-to-commission line within the right-of-way."),
 scope=["The contractor shall clear and grade the right-of-way and establish access.",
   "The contractor shall string, bend, and weld pipe to the referenced welding procedure specification.",
   "The contractor shall perform non-destructive examination of welds per the project specification.",
   "The contractor shall apply field joint coating and inspect for holidays before lowering-in.",
   "The contractor shall trench, lower-in, and backfill to the specified depth of cover.",
   "The contractor shall install pig launchers/receivers and tie-ins as shown on the drawings.",
   "The contractor shall hydrotest the segment and provide signed test records before commissioning."],
 exclusions=["Engineering and design, and permitting, are by others unless stated.",
   "Material procurement of line pipe is by SOWsmith unless the bid documents state otherwise.",
   "Final commissioning and gas-in is excluded from this scope."],
 deliverables=("The contractor shall deliver a complete, hydrotested pipeline segment with as-built "
   "drawings, weld maps, NDE reports, coating inspection records, and hydrotest documentation per the "
   "schedule below."),
 schedule_table=("Indicative milestones (durations confirmed at award):",
   ["Milestone", "Deliverable", "Timing"],
   [["Mobilization", "Crew and equipment on site", "Day 1"],
    ["Welding complete", "Weld map + NDE reports", "[TO CONFIRM]"],
    ["Lower-in & backfill", "Depth-of-cover survey", "[TO CONFIRM]"],
    ["Hydrotest", "Signed test records", "[TO CONFIRM]"]]),
 materials=("SOWsmith supplies line pipe, valves and major fittings unless the bid documents state "
   "otherwise. The contractor supplies consumables, welding equipment, and construction spread. "
   "Site access, laydown, and operating conditions are as described in the project documents."),
 safety=_SAFETY,
 acceptance=("Work is accepted upon successful hydrotest, completion of punch-list items, and submission "
   "of the complete records package. Acceptance follows the referenced standards:"),
 references=[f"{ORG} Welding Services — Statement of Work", f"{ORG} Hydrotesting & Commissioning — Statement of Work",
   f"{ORG} Standard Bid Clause Library", f"{ORG} Bidder Instructions"])

# -------- EXEMPLAR 2: Facility Maintenance ---------------------------------
add(type="Facility Maintenance", exemplar=True, status=STANDARD, doc="SOW",
 file="Facility-Maintenance-SOW", title=f"{ORG} Facility Maintenance — Statement of Work (Exemplar)",
 docid="BP-SOW-FM-002", version="v1.3", date="2025-11-03",
 overview=("Planned and corrective maintenance at an operating compressor or processing facility during "
   "a scheduled turnaround window, restoring equipment to service safely and on schedule."),
 scope=["The contractor shall service valves, actuators, and rotating equipment per the work list.",
   "The contractor shall perform minor mechanical repairs and component replacements as directed.",
   "The contractor shall surface-prep and paint designated equipment and structures.",
   "The contractor shall support lockout/tagout and confined-space entries per site procedures.",
   "The contractor shall clean up and restore work areas at the end of each shift."],
 exclusions=["Major overhauls and capital replacements are excluded unless added by change order.",
   "Engineering, spare-parts procurement, and electrical isolation are by others unless stated."],
 deliverables=("The contractor shall deliver completed work-list items, red-lined records of repairs, "
   "and a turnaround close-out report within the maintenance window."),
 materials=("SOWsmith provides facility access, isolation, and major spares; the contractor provides labor, "
   "tools, and consumables. Work occurs in an operating facility with strict permit controls."),
 safety=_SAFETY,
 acceptance=("Work is accepted on completion of the work list, return-to-service checks, and the close-out "
   "report. Acceptance follows the referenced standards:"),
 references=[f"{ORG} Standard Bid Clause Library", f"{ORG} Bidder Instructions"])

# -------- EXEMPLAR 3: Electrical & Instrumentation -------------------------
add(type="Electrical & Instrumentation", exemplar=True, status=STANDARD, doc="SOW",
 file="Electrical-Instrumentation-SOW", title=f"{ORG} Electrical & Instrumentation — Statement of Work (Exemplar)",
 docid="BP-SOW-EI-003", version="v1.2", date="2025-09-20",
 overview=("Installation and termination of electrical and instrumentation systems for a new pump or "
   "metering station, delivered tested and ready for loop checks and commissioning support."),
 scope=["The contractor shall install cable tray, conduit, and supports per the drawings.",
   "The contractor shall pull, terminate, and label power, control, and instrument cabling.",
   "The contractor shall install field instruments and junction boxes per the specifications.",
   "The contractor shall perform continuity, megger, and loop checks and document results.",
   "The contractor shall support commissioning and functional testing as directed."],
 exclusions=["Control-system configuration and PLC programming are by others unless stated.",
   "Instrument supply is by SOWsmith unless the bid documents state otherwise."],
 deliverables=("The contractor shall deliver installed and tested electrical and instrumentation systems, "
   "loop-check sheets, megger records, and red-lined drawings."),
 materials=("SOWsmith supplies major instruments and panels; the contractor supplies cable, tray, conduit, "
   "fittings, and labor. Classified-area requirements apply where indicated."),
 safety=_SAFETY,
 acceptance=("Work is accepted on successful loop checks, documentation submission, and energization "
   "readiness. Acceptance follows the referenced standards:"),
 references=[f"{ORG} Standard Bid Clause Library", f"{ORG} Bidder Instructions"])

# -------- EXEMPLAR 4: Civil & Earthwork ------------------------------------
add(type="Civil & Earthwork", exemplar=True, status=STANDARD, doc="SOW",
 file="Civil-Earthwork-SOW", title=f"{ORG} Civil & Earthwork — Statement of Work (Exemplar)",
 docid="BP-SOW-CE-004", version="v2.0", date="2026-02-10",
 overview=("Site preparation, grading, access road, and foundation earthwork for a new facility pad, "
   "delivered to line and grade with documented compaction and erosion control in place."),
 scope=["The contractor shall clear, strip, and grade the site to the specified line and grade.",
   "The contractor shall construct the access road and laydown areas per the drawings.",
   "The contractor shall excavate and prepare foundation subgrades and place structural fill.",
   "The contractor shall install erosion and sediment controls and maintain them.",
   "The contractor shall provide compaction testing records for all placed fill."],
 exclusions=["Concrete foundations, structures, and surveying control are by others unless stated.",
   "Permitting and geotechnical investigation are by SOWsmith unless the bid documents state otherwise."],
 deliverables=("The contractor shall deliver a graded, compacted site with as-built grades, compaction "
   "test reports, and erosion-control documentation."),
 materials=("SOWsmith provides survey control and import-fill source where specified; the contractor "
   "provides equipment, labor, and erosion-control materials. Weather and access conditions apply."),
 safety=_SAFETY,
 acceptance=("Work is accepted on grade verification, passing compaction tests, and documentation "
   "submission. Acceptance follows the referenced standards:"),
 references=[f"{ORG} Standard Bid Clause Library", f"{ORG} Bidder Instructions"])

# -------- EXEMPLAR 5: Coating & Insulation ---------------------------------
add(type="Coating & Insulation", exemplar=True, status=STANDARD, doc="SOW",
 file="Coating-Insulation-SOW", title=f"{ORG} Coating & Insulation — Statement of Work (Exemplar)",
 docid="BP-SOW-CI-005", version="v1.1", date="2025-12-05",
 overview=("Surface preparation, protective coating, and insulation of piping and equipment to the "
   "specified system, delivered with documented inspection of every coated surface."),
 scope=["The contractor shall surface-prep to the specified standard (e.g., abrasive blast profile).",
   "The contractor shall apply the specified coating system to the required dry film thickness.",
   "The contractor shall measure and record dry film thickness and environmental conditions.",
   "The contractor shall install insulation and weatherproofing where specified.",
   "The contractor shall perform holiday detection and document all inspection results."],
 exclusions=["Coating material is by SOWsmith unless the bid documents state otherwise.",
   "Scaffolding and access may be by others where indicated."],
 deliverables=("The contractor shall deliver coated and insulated systems with dry-film-thickness logs, "
   "environmental records, and holiday-detection results."),
 materials=("SOWsmith supplies coating and insulation materials where specified; the contractor supplies "
   "blast media, equipment, and labor. Ambient/weather windows govern application."),
 safety=_SAFETY,
 acceptance=("Work is accepted on passing inspection (thickness, adhesion, holiday) and documentation "
   "submission. Acceptance follows the referenced standards:"),
 references=[f"{ORG} Standard Bid Clause Library", f"{ORG} Bidder Instructions"])

# ============================ REFERENCE-ONLY DOCS ==========================
def _ref(stype, file, title, docid, version, date, overview, scope, doc="SOW", status=STANDARD,
         exclusions=None, special=None):
    add(type=stype, exemplar=False, status=status, doc=doc, file=file, title=title, docid=docid,
        version=version, date=date, overview=overview, scope=scope,
        exclusions=exclusions or ["Work not expressly listed is excluded.",
                                  "Engineering, permitting, and material supply are by others unless stated."],
        deliverables=("The contractor shall deliver the completed work with the records and as-builts "
                      "required by the referenced specifications." + (f" {special}" if special else "")),
        materials=("Material and equipment supply is split between SOWsmith and the contractor as stated "
                   "in the bid documents; site and access conditions are as provided."),
        safety=_SAFETY,
        acceptance="Work is accepted on completion, inspection, and documentation per the referenced standards:",
        references=[f"{ORG} Standard Bid Clause Library", f"{ORG} Bidder Instructions"])

_ref("Hydrotesting & Commissioning", "Hydrotesting-Commissioning-SOW",
     f"{ORG} Hydrotesting & Commissioning — Statement of Work", "BP-SOW-HC-006", "v1.4", "2025-07-18",
     "Pressure testing and commissioning of a pipeline or facility system to confirm integrity and readiness.",
     ["The contractor shall fill, pressurize, and hold the system to the specified test pressure.",
      "The contractor shall monitor and record pressure and temperature for the hold period.",
      "The contractor shall dewater, dry, and clean the system after a successful test.",
      "The contractor shall provide signed test and commissioning records."])
_ref("Fabrication", "Fabrication-SOW", f"{ORG} Shop Fabrication — Statement of Work",
     "BP-SOW-FAB-007", "v1.0", "2026-01-09",
     "Shop fabrication of pipe spools and skid assemblies to the issued isometrics and specifications.",
     ["The contractor shall fabricate spools and assemblies per the isometrics and welding procedures.",
      "The contractor shall perform non-destructive examination and dimensional checks.",
      "The contractor shall coat and prepare assemblies for shipment.",
      "The contractor shall provide material traceability and fabrication records."])
_ref("Cathodic Protection", "Cathodic-Protection-SOW", f"{ORG} Cathodic Protection — Statement of Work",
     "BP-SOW-CP-008", "v1.2", "2025-10-12",
     "Installation and testing of cathodic protection systems to protect buried or submerged assets.",
     ["The contractor shall install anodes, rectifiers, and test stations per the drawings.",
      "The contractor shall bond and connect the system and verify continuity.",
      "The contractor shall take baseline potential readings and document results.",
      "The contractor shall provide commissioning and survey records."])
_ref("Pipeline Integrity", "Pipeline-Integrity-Inspection-SOW",
     f"{ORG} Pipeline Integrity & Inspection — Statement of Work", "BP-SOW-PI-009", "v1.1", "2025-08-15",
     "In-line and field inspection (NDE) to assess pipeline integrity and support fitness-for-service.",
     ["The contractor shall perform the specified inspection and non-destructive examination methods.",
      "The contractor shall locate, characterize, and report anomalies per the procedure.",
      "The contractor shall provide calibrated equipment and qualified technicians.",
      "The contractor shall deliver an inspection report with findings and data."])
_ref("Demolition & Abandonment", "Demolition-Abandonment-SOW",
     f"{ORG} Demolition & Abandonment — Statement of Work", "BP-SOW-DA-010", "v1.0", "2025-12-01",
     "Safe demolition, removal, and abandonment-in-place of pipeline or facility assets per regulations.",
     ["The contractor shall isolate, purge, and make safe all systems before work begins.",
      "The contractor shall demolish, cut, and remove assets and dispose of materials lawfully.",
      "The contractor shall abandon-in-place per the approved procedure where specified.",
      "The contractor shall restore the site and provide completion records."])
_ref("Station & Terminal", "Station-Terminal-Construction-SOW",
     f"{ORG} Station & Terminal Construction — Statement of Work", "BP-SOW-ST-011", "v1.5", "2025-11-28",
     "Construction of station/terminal mechanical, piping, and structural scope to deliver a complete facility.",
     ["The contractor shall install piping, supports, and equipment per the drawings.",
      "The contractor shall erect structural steel and platforms as shown.",
      "The contractor shall support testing, flushing, and commissioning activities.",
      "The contractor shall provide as-builts and turnover documentation."])
_ref("Welding Services", "Welding-Services-SOW", f"{ORG} Welding Services — Statement of Work",
     "BP-SOW-WS-012", "v2.3", "2025-09-05",
     "Provision of qualified welding services to the applicable welding procedure specifications.",
     ["The contractor shall provide qualified welders and welding procedure specifications.",
      "The contractor shall perform production welds and support non-destructive examination.",
      "The contractor shall maintain weld logs and welder qualification records.",
      "The contractor shall meet the acceptance criteria of the referenced code."])

# RFP example with non-standard commercial terms -> LEGAL REVIEW REQUIRED
add(type="Facility Maintenance", exemplar=False, status=LEGAL, doc="RFP",
 file="Facility-Maintenance-RFP", title=f"{ORG} Facility Maintenance — Request for Proposal (Example)",
 docid="BP-RFP-FM-013", version="v1.0", date="2026-01-20",
 overview=("Request for Proposal for multi-year facility maintenance services, including instructions to "
   "bidders, evaluation criteria, and the commercial terms bidders must accept."),
 scope=["Bidders shall propose to perform the facility maintenance scope described in the attached SOW.",
   "Bidders shall provide a performance bond and meet the stated insurance limits. [TERMS - LEGAL REVIEW]",
   "Bidders shall accept the warranty and indemnification terms in the contract. [TERMS - LEGAL REVIEW]",
   "Bidders shall submit pricing on the provided rate-sheet template.",
   "Bidders shall describe safety performance, references, and capacity."],
 exclusions=["Scope not in the attached SOW is excluded.",
   "Specific bonding amounts, insurance limits, and liability caps are subject to Legal/Contracts review."],
 deliverables=("Bidders shall deliver a complete proposal: technical approach, completed rate sheet, "
   "bonding/insurance evidence, safety record, and signed acknowledgment of terms."),
 materials=("As described in the attached SOW. Commercial and legal terms are governed by the contract "
   "and are subject to Legal/Contracts review before issue."),
 safety=_SAFETY,
 acceptance=("Proposals are evaluated on price, technical merit, safety, and acceptance of terms. Award is "
   "subject to Legal/Contracts approval. Acceptance follows the referenced standards:"),
 references=[f"{ORG} Facility Maintenance — Statement of Work", f"{ORG} Standard Bid Clause Library",
   f"{ORG} Bidder Instructions"])

# ---------------------------------------------------------------- write docs
def reset_dir():
    if os.path.exists(ROOT): shutil.rmtree(ROOT)
    for d in ["Approved-Exemplars", "Reference-Library", "Templates", "Lists",
              "Drafts", "Approved", "Rejected"]:
        os.makedirs(os.path.join(ROOT, d), exist_ok=True)

def fname(spec):
    return f"{spec['file']}_{spec['version']}_{spec['date']}.docx"

def write_docs():
    for spec in D:
        f = fname(spec)
        ref = os.path.join(ROOT, "Reference-Library", f)
        build_sow(spec, ref)
        if spec.get("exemplar"):
            shutil.copy2(ref, os.path.join(ROOT, "Approved-Exemplars", f))

# ---------------------------------------------------------------- templates + clause library
def write_template():
    doc = new_doc()
    title_block(doc, f"{ORG} Bid Package (SOW/RFP) Template", f"{ORG} · House template · Supply Chain")
    control_table(doc, [("Document ID", "BP-TMPL-SOW-001"), ("Version", "v1.0"),
        ("Effective date", "2026-01-01"), ("Owner", "Supply Chain"),
        ("Purpose", "Canonical structure every bid package must follow")])
    para(doc, "Keep the seven headings, in this order, and the version footer exactly as shown so it "
              "can be parsed by the approval workflow. Replace the guidance with real content.")
    guide = {
      "1. PROJECT OVERVIEW": "What the project is, where, and the objective.",
      "2. SCOPE OF WORK": "Numbered work items. Specific and unambiguous.",
      "3. EXCLUSIONS & ASSUMPTIONS": "What is NOT in scope and the assumptions. Never leave empty.",
      "4. DELIVERABLES & SCHEDULE": "Deliverables, milestones, durations. Use placeholders for missing dates.",
      "5. MATERIALS, EQUIPMENT & SITE CONDITIONS": "Who supplies what; site/access conditions.",
      "6. SAFETY, QUALITY & COMPLIANCE": "Safety-first requirements, quality standards, OSHA/DOT/PHMSA.",
      "7. ACCEPTANCE & REFERENCES": "Acceptance criteria and related documents (use exact titles).",
    }
    for h, g in guide.items():
        heading(doc, h); para(doc, g)
    footer_block(doc, "v0.1 DRAFT", "YYYY-MM-DD", "[service_type]", "[SOW | RFP]",
                 "[STANDARD REVIEW | LEGAL REVIEW REQUIRED]")
    doc.save(os.path.join(ROOT, "Templates", "_TEMPLATE_Bid-Package-SOW_v1.0.docx"))

def write_clause_library():
    doc = new_doc()
    title_block(doc, f"{ORG} Standard Bid Clause Library", f"{ORG} · Reference · Supply Chain / Legal")
    control_table(doc, [("Document ID", "BP-REF-CLAUSE-001"), ("Version", "v1.0"),
        ("Effective date", "2026-01-01"), ("Owner", "Legal / Contracts")])
    para(doc, "Approved, reusable clause language for bid packages. Use STANDARD clauses freely. Any "
              "NON-STANDARD clause (below) requires Legal/Contracts review before issue.")
    heading(doc, "Standard clauses")
    bullets(doc, ["Scope changes are handled by written change order only.",
        "The contractor follows SOWsmith's safety-first program and site permit-to-work controls.",
        "Work is performed to the referenced specifications and inspection/test plan.",
        "Invoicing follows the agreed schedule of values; standard payment terms are Net 30."])
    heading(doc, "Non-standard clauses — LEGAL / CONTRACTS REVIEW REQUIRED")
    bullets(doc, ["Liquidated damages or schedule penalties.",
        "Indemnification, hold-harmless, or liability caps/waivers.",
        "Performance/payment bonds, insurance limits, retainage, or lien terms.",
        "Non-standard payment terms (e.g., Net 60/90), termination-for-convenience, unusual flow-downs."])
    footer_block(doc, "v1.0", "2026-01-01", "All", "Reference", "STANDARD REVIEW")
    doc.save(os.path.join(ROOT, "Templates", "Standard-Bid-Clause-Library_v1.0.docx"))

def write_bidder_form():
    doc = new_doc()
    title_block(doc, f"{ORG} Bidder Response Form", f"{ORG} · Form · Supply Chain")
    para(doc, "Bidders complete and return this form with their proposal.")
    t = doc.add_table(rows=0, cols=2); t.style = "Table Grid"
    for label in ["Bidder company", "Contact name / email", "Service type",
                  "Total proposed price", "Schedule / duration offered", "Exclusions / clarifications",
                  "Bonding & insurance (if required)", "Safety record (EMR)", "References",
                  "Acknowledgment of terms (Y/N)", "Authorized signature / date"]:
        c = t.add_row().cells; rr = c[0].paragraphs[0].add_run(label); rr.bold = True; rr.font.size = Pt(10)
        _shade(c[0], "F1F5F7"); c[1].paragraphs[0].add_run(" ")
    footer_block(doc, "v1.0", "2026-01-01", "All", "Form", "STANDARD REVIEW")
    doc.save(os.path.join(ROOT, "Templates", "Bidder-Response-Form_v1.0.docx"))

# ---------------------------------------------------------------- multi-format reference resources
MEDIA = []
def note_media(rel, kind, desc): MEDIA.append((rel, kind, desc))

def write_bidder_instructions_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    path = os.path.join(ROOT, "Reference-Library", "Bidder-Instructions_v1.0_2026-01-01.pdf")
    ss = getSampleStyleSheet()
    h = ParagraphStyle("h", parent=ss["Title"], textColor=HexColor("#0C3B4D"), fontSize=20)
    h2 = ParagraphStyle("h2", parent=ss["Heading2"], textColor=HexColor("#0C3B4D"))
    body = ParagraphStyle("b", parent=ss["BodyText"], fontSize=10.5, leading=15)
    el = [Paragraph(f"{ORG} Bidder Instructions", h), Spacer(1, 4),
          Paragraph("How to respond to a SOWsmith bid package. For full scope, always refer to the "
                    "specific Statement of Work (SOW) or Request for Proposal (RFP) issued to you.", body),
          Spacer(1, 8)]
    for t, txt in [
        ("Submitting a bid", "Complete the Bidder Response Form, attach pricing on the provided rate sheet, "
         "and acknowledge the terms. Late submissions may be rejected."),
        ("Scope & exclusions", "Bid exactly the Scope of Work. State exclusions and assumptions clearly — "
         "the Exclusions & Assumptions section governs disputes."),
        ("Safety first", "SOWsmith awards work to bidders who put safety first. Provide your safety record "
         "and confirm compliance with OSHA, DOT, and PHMSA requirements."),
        ("Commercial & legal terms", "Standard terms are in the Standard Bid Clause Library. Any "
         "non-standard term (bonding, indemnification, liquidated damages) is subject to Legal/Contracts review."),
        ("Award", "Award is based on price, technical merit, safety, and acceptance of terms, and is "
         "approved by Supply Chain (and Legal/Contracts where terms are non-standard).")]:
        el += [Paragraph(t, h2), Paragraph(txt, body), Spacer(1, 6)]
    el += [Spacer(1, 8), Paragraph("Document ID BP-REF-BID-001 · v1.0 · 2026-01-01 · Owner: Supply Chain · "
           "Classification: Internal", ParagraphStyle("f", parent=body, fontSize=8.5, textColor=HexColor("#5B727C")))]
    SimpleDocTemplate(path, pagesize=A4, leftMargin=22*mm, rightMargin=22*mm, topMargin=20*mm,
                      bottomMargin=18*mm, title=f"{ORG} Bidder Instructions").build(el)
    note_media("Reference-Library/Bidder-Instructions_v1.0_2026-01-01.pdf", "PDF",
               "How bidders respond to a SOWsmith bid package")

NAVY_P, AMBER_P, GREY_P, INK_P, MIST_P = (PColor(0x0C,0x3B,0x4D), PColor(0xE0,0x7B,0x27),
    PColor(0x5B,0x72,0x7C), PColor(0x1F,0x2D,0x33), PColor(0xCF,0xDE,0xE4))

def write_bid_process_pptx():
    prs = Presentation(); prs.slide_width = PIn(13.333); prs.slide_height = PIn(7.5)
    def cover(kick, title, sub):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY_P; bar.line.fill.background(); bar.shadow.inherit = False
        tf = s.shapes.add_textbox(PIn(0.8), PIn(2.6), PIn(11.7), PIn(3)).text_frame; tf.word_wrap = True
        for i,(t,sz,c,b) in enumerate([(kick,15,AMBER_P,True),(title,40,PColor(0xFF,0xFF,0xFF),True),(sub,18,MIST_P,False)]):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); r = p.add_run(); r.text=t; r.font.size=PPt(sz); r.font.bold=b; r.font.color.rgb=c; p.space_after=PPt(8)
    def bullets_slide(title, items):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        th = s.shapes.add_textbox(PIn(0.7), PIn(0.45), PIn(12), PIn(1)).text_frame; tr = th.paragraphs[0].add_run()
        tr.text=title; tr.font.size=PPt(28); tr.font.bold=True; tr.font.color.rgb=NAVY_P
        bf = s.shapes.add_textbox(PIn(0.75), PIn(1.7), PIn(11.9), PIn(5)).text_frame; bf.word_wrap=True
        for i,it in enumerate(items):
            p = bf.paragraphs[0] if i==0 else bf.add_paragraph(); r=p.add_run(); r.text="•  "+it; r.font.size=PPt(17); r.font.color.rgb=INK_P; p.space_after=PPt(7)
    cover(f"{ORG} · SUPPLY CHAIN", "Bid Package Process", "From project notes to issued SOW/RFP")
    bullets_slide("The bid package flow", [
        "A buyer submits project notes (service type, scope, schedule).",
        "The Bid Package Agent drafts a house-style SOW/RFP grounded in approved exemplars.",
        "An ambiguity check flags vague language, missing exclusions, and undefined terms.",
        "Supply Chain reviews and approves; Legal/Contracts reviews any non-standard terms.",
        "The approved package is issued to bidders; every action is logged."])
    bullets_slide("What to get right", [
        "Scope and exclusions must be specific — they prevent disputes.",
        "Never invent prices or commercial terms; leave placeholders for Estimating/Contracts.",
        "Safety-first requirements appear in every package.",
        "Reuse approved clause language; flag non-standard clauses for Legal."])
    rel = "Reference-Library/Bid-Process-Overview_v1.0_2026-01-02.pptx"
    prs.save(os.path.join(ROOT, rel))
    note_media(rel, "PPTX", "Overview of the bid-package process for buyers")

def write_rate_sheet_xlsx():
    wb = Workbook(); ws = wb.active; ws.title = "Rate Sheet (template)"
    ws.append(["Item", "Description", "Unit", "Est. Qty", "Unit Price (USD)", "Extended (USD)"])
    style_header(ws, 6)
    for r in [["1", "Mobilization / demobilization", "LS", 1, "", ""],
              ["2", "Pipe welding", "weld", "[QTY]", "", ""],
              ["3", "Trench / lower-in / backfill", "LF", "[QTY]", "", ""],
              ["4", "Hydrotest", "test", 1, "", ""]]:
        ws.append(r)
    ws.append([]); ws.append(["NOTE", "Bidders enter unit prices. SOWsmith does not pre-fill pricing.", "", "", "", ""])
    for col, w in zip("ABCDEF", [8, 36, 8, 12, 16, 16]):
        ws.column_dimensions[col].width = w
    ws.freeze_panes = "A2"
    wb.save(os.path.join(ROOT, "Reference-Library", "Rate-Sheet-Template_v1.0.xlsx"))
    note_media("Reference-Library/Rate-Sheet-Template_v1.0.xlsx", "XLSX",
               "Blank bidder rate sheet (pricing entered by bidders, not by SOWsmith)")

# ---------------------------------------------------------------- lists (xlsx)
# service_type -> (Supply Chain reviewer, Legal/Contracts reviewer)
REVIEWERS = [
  ("Pipeline Construction", "mark.davis@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Facility Maintenance", "jen.alvarez@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Electrical & Instrumentation", "tom.nguyen@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Civil & Earthwork", "mark.davis@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Coating & Insulation", "jen.alvarez@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Hydrotesting & Commissioning", "tom.nguyen@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Fabrication", "mark.davis@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Cathodic Protection", "tom.nguyen@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Pipeline Integrity", "jen.alvarez@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Demolition & Abandonment", "mark.davis@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Station & Terminal", "jen.alvarez@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
  ("Welding Services", "tom.nguyen@sowsmithusa.com", "robert.hayes@sowsmithusa.com"),
]

AUDIT_SAMPLE = [
  ("Pipeline-Construction_SOW_v0.1_DRAFT_2026-05-12.md", "Mark Davis", "mark.davis@sowsmithusa.com",
   "Approved", "2026-05-12T14:32:00Z", "STANDARD REVIEW", "Scope clear; approved for issue."),
  ("Facility-Maintenance_RFP_v0.1_DRAFT_2026-05-13.md", "Jen Alvarez", "robert.hayes@sowsmithusa.com",
   "Request Changes", "2026-05-13T09:10:00Z", "LEGAL REVIEW REQUIRED", "Bonding/insurance terms to Legal."),
  ("Civil-Earthwork_SOW_v0.2_DRAFT_2026-05-15.md", "Mark Davis", "mark.davis@sowsmithusa.com",
   "Approved", "2026-05-15T11:05:00Z", "STANDARD REVIEW", "Compaction testing added; approved."),
  ("Pipeline-Construction_SOW_v0.1_DRAFT_2026-05-16.md", "Tom Nguyen", "robert.hayes@sowsmithusa.com",
   "Rejected", "2026-05-16T13:20:00Z", "LEGAL REVIEW REQUIRED", "Liquidated-damages clause not approved; redraft."),
  ("(system)", "(flow)", "(flow)", "ERROR", "2026-05-16T13:21:30Z", "STANDARD REVIEW",
   "SharePoint move timed out; retried 3x then alerted owner."),
]

def style_header(ws, ncols):
    fill = PatternFill("solid", fgColor="0C3B4D"); side = Side(style="thin", color="B7C4CB")
    for c in range(1, ncols + 1):
        cell = ws.cell(row=1, column=c); cell.font = Font(bold=True, color="FFFFFF", size=11)
        cell.fill = fill; cell.alignment = Alignment(vertical="center")
        cell.border = Border(left=side, right=side, top=side, bottom=side)
    ws.row_dimensions[1].height = 22

def write_reviewers_xlsx():
    wb = Workbook(); ws = wb.active; ws.title = "Reviewers"
    ws.append(["ServiceType", "ReviewerEmail", "LegalReviewerEmail"]); style_header(ws, 3)
    for row in REVIEWERS: ws.append(list(row))
    ws.column_dimensions["A"].width = 30; ws.column_dimensions["B"].width = 34; ws.column_dimensions["C"].width = 34
    ws.freeze_panes = "A2"; wb.save(os.path.join(ROOT, "Lists", "Reviewers.xlsx"))

def write_audit_xlsx():
    wb = Workbook(); ws = wb.active; ws.title = "AuditLog"
    ws.append(["DraftFileName", "Author", "Reviewer", "Action", "Timestamp", "ReviewStatus", "Comments"])
    style_header(ws, 7)
    for row in AUDIT_SAMPLE: ws.append(list(row))
    for i, w in enumerate([46, 16, 34, 16, 22, 22, 46], start=1):
        ws.column_dimensions[chr(64 + i)].width = w
    ws.freeze_panes = "A2"; wb.save(os.path.join(ROOT, "Lists", "AuditLog.xlsx"))

# ---------------------------------------------------------------- output libs + readme
def write_output_readmes():
    notes = {
      "Drafts": ("Drafts library (output store).\n\nThe agent writes new bid-package drafts here. "
        "Creating a file here TRIGGERS the Power Automate approval flow. Starts empty.\nNot a grounding source."),
      "Approved": ("Approved library (issued bid packages).\n\nPower Automate moves a draft here on approval "
        "and increments the version. Starts empty.\nNot a grounding source."),
      "Rejected": ("Rejected library.\n\nPower Automate moves a draft here on rejection; the author may "
        "resubmit. Starts empty.\nNot a grounding source."),
    }
    for lib, txt in notes.items():
        with open(os.path.join(ROOT, lib, "_README.txt"), "w", encoding="utf-8") as f:
            f.write(txt + "\n")

def write_root_readme():
    ex = [s for s in D if s.get("exemplar")]
    rows_ex = "\n".join(f"| {s['type']} | {s['title']} | {s.get('doc','SOW')} | {s['version']} | `{fname(s)}` |" for s in ex)
    rows_ref = "\n".join(f"| {s['type']} | {s['title']} | {s.get('doc','SOW')} | {s['status']} | `{fname(s)}` |" for s in D)
    rows_media = "\n".join(f"| {kind} | `{rel.split('/')[-1]}` | {desc} |" for rel, kind, desc in MEDIA)
    md = f"""# SharePoint Online Knowledge Base — Bid Package Generator (SOW/RFP)

This folder is the **knowledge base** the Bid Package Agent grounds on. Upload it to a SharePoint
Online site (recommended site name **`POC-BidPackage`**) and wire each library as described below.

> **Why two grounding libraries?** `Approved-Exemplars` ({len(ex)} gold-standard prior SOWs) is what
> the **generator (Flow B)** copies house style and structure from. `Reference-Library` (the broad set
> of prior bid documents + the clause library) is what **Bid-document Q&A (Flow A)** grounds in. Keep
> them separate — mixing them dilutes draft quality.

## Folder map

| Folder | Role | Wire it to |
|---|---|---|
| `Approved-Exemplars/` | {len(ex)} gold-standard SOWs — the generator's style + factual anchor | Topic 1 *Draft a Bid Package* → Generative Answers knowledge (this library **only**) |
| `Reference-Library/` | broad corpus ({len(D)} bid docs + clause library + resources) | Topic 2 *Bid-document Q&A* → knowledge source |
| `Templates/` | SOW/RFP template, Standard Bid Clause Library, Bidder Response Form | "Get official templates" path |
| `Lists/` | `Reviewers.xlsx`, `AuditLog.xlsx` — import as **SharePoint lists** | Power Automate actions 4 & 10 |
| `Drafts/` `Approved/` `Rejected/` | output libraries — start empty | Power Automate trigger + destinations |

## `service_type` tokens (must match exactly)
Every document's footer `SERVICE TYPE:` value matches a row in `Lists/Reviewers.xlsx` (`ServiceType`).
Power Automate looks up the reviewer with `ServiceType eq '<token>'`. Tokens:
{", ".join(sorted(set(s['type'] for s in D)))}.

> **Review tier is content-driven:** a package is flagged **LEGAL REVIEW REQUIRED** when it contains
> non-standard commercial/legal terms (liquidated damages, indemnification, bonding, retainage,
> non-standard payment terms) — routed to Legal/Contracts. Otherwise **STANDARD REVIEW** (Supply Chain).

## Approved-Exemplars (feeds the generator — Flow B)

| Service type | Title | Doc | Version | File |
|---|---|---|---|---|
{rows_ex}

## Reference-Library (feeds Bid-document Q&A — Flow A)

| Service type | Title | Doc | Review status | File |
|---|---|---|---|---|
{rows_ref}

## Multi-format resources

| Type | File | What it is |
|---|---|---|
{rows_media}

## Upload checklist (SharePoint Online + Copilot Studio)
1. Create the site `POC-BidPackage` (default-themed; custom CSS / accordion nav can break grounding).
2. Create document libraries `Approved-Exemplars`, `Reference-Library`, `Templates`, `Drafts`,
   `Approved`, `Rejected` and upload the matching folder contents.
3. Create the lists `Reviewers` and `AuditLog` (Lists → "From Excel"). `ServiceType` = Choice.
4. In Copilot Studio, add SharePoint knowledge using the **bare URL** (no `https://`). Scope Topic 1 to
   `Approved-Exemplars` only and Topic 2 to `Reference-Library`.
5. **File-size limit:** grounding files must be **< 7 MB** without an M365 Copilot license (< 200 MB with
   Copilot + Work IQ). These files are tiny; verify the tenant setting with the admin.

_Generated by `_build_knowledge_base.py`. Re-run to regenerate. Organization: {ORG}. US English._
"""
    with open(os.path.join(ROOT, "README.md"), "w", encoding="utf-8") as f:
        f.write(md)

# ---------------------------------------------------------------- main
def main():
    MEDIA.clear()
    reset_dir()
    write_docs()
    write_template(); write_clause_library(); write_bidder_form()
    write_bidder_instructions_pdf(); write_bid_process_pptx(); write_rate_sheet_xlsx()
    write_reviewers_xlsx(); write_audit_xlsx()
    write_output_readmes(); write_root_readme()
    from collections import Counter
    kinds = Counter(os.path.splitext(f)[1].lower() for dp, _, fs in os.walk(ROOT) for f in fs)
    print(f"OK — built bid-package knowledge base at: {ROOT}")
    print(f"   Approved-Exemplars: {len([s for s in D if s.get('exemplar')])} SOWs (DOCX)")
    print(f"   Reference-Library : {len(D)} bid docs + {len(MEDIA)} resources")
    print(f"   File types        : " + ", ".join(f"{k or 'none'}:{v}" for k, v in sorted(kinds.items())))
    print(f"   Total files       : {sum(kinds.values())}")

if __name__ == "__main__":
    main()
